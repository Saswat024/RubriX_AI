"""
Document parser for extracting text from various file formats
Supports: .docx, .pptx, .pdf, .txt
Also supports extracting images from PDFs for flowchart processing
"""

import io
import base64
from typing import Dict, Any, Tuple, Optional
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader
import fitz  # pymupdf for image extraction


async def parse_document(file_base64: str, file_type: str) -> Dict[str, Any]:
    """
    Parse document and extract text content

    Args:
        file_base64: Base64 encoded file content
        file_type: File extension (.docx, .pptx, .pdf, .txt)

    Returns:
        Dictionary containing extracted text and metadata
    """
    try:
        # Decode base64 content
        file_content = base64.b64decode(
            file_base64.split(",")[1] if "," in file_base64 else file_base64
        )
        file_stream = io.BytesIO(file_content)

        extracted_text = ""
        metadata = {"type": file_type, "pages": 0, "slides": 0, "paragraphs": 0}

        if file_type == ".docx":
            extracted_text, metadata = await parse_docx(file_stream)
        elif file_type == ".pptx":
            extracted_text, metadata = await parse_pptx(file_stream)
        elif file_type == ".pdf":
            extracted_text, metadata = await parse_pdf(file_stream)
        elif file_type == ".txt":
            extracted_text = file_content.decode("utf-8")
            metadata["paragraphs"] = len(extracted_text.split("\n\n"))
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

        return {"text": extracted_text, "metadata": metadata, "success": True}

    except Exception as e:
        print(f"Error parsing document: {str(e)}")
        return {"text": "", "metadata": {}, "success": False, "error": str(e)}


async def parse_docx(file_stream: io.BytesIO) -> Tuple[str, Dict[str, Any]]:
    """Parse DOCX file and extract text"""
    doc = Document(file_stream)
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
    text = "\n\n".join(paragraphs)

    metadata = {
        "type": ".docx",
        "paragraphs": len(paragraphs),
        "pages": len(doc.sections),
    }

    return text, metadata


async def parse_pptx(file_stream: io.BytesIO) -> Tuple[str, Dict[str, Any]]:
    """Parse PPTX file and extract text from slides"""
    prs = Presentation(file_stream)
    slide_texts = []

    for i, slide in enumerate(prs.slides):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text)
        if slide_content:
            slide_texts.append(f"Slide {i + 1}:\n" + "\n".join(slide_content))

    text = "\n\n".join(slide_texts)

    metadata = {"type": ".pptx", "slides": len(prs.slides), "pages": len(prs.slides)}

    return text, metadata


async def parse_pdf(file_stream: io.BytesIO) -> Tuple[str, Dict[str, Any]]:
    """Parse PDF file and extract text"""
    pdf_reader = PdfReader(file_stream)
    page_texts = []

    for i, page in enumerate(pdf_reader.pages):
        page_text = page.extract_text()
        if page_text.strip():
            page_texts.append(f"Page {i + 1}:\n{page_text}")

    text = "\n\n".join(page_texts)

    metadata = {"type": ".pdf", "pages": len(pdf_reader.pages)}

    return text, metadata


async def parse_pdf_for_images(file_base64: str) -> Dict[str, Any]:
    """
    Extract images from PDF file for flowchart processing.
    Returns the first/largest image found as base64.

    Args:
        file_base64: Base64 encoded PDF file

    Returns:
        Dictionary with image data or error
    """
    try:
        # Decode base64 content
        file_content = base64.b64decode(
            file_base64.split(",")[1] if "," in file_base64 else file_base64
        )

        # Open PDF with pymupdf
        pdf_document = fitz.open(stream=file_content, filetype="pdf")

        best_image = None
        best_image_size = 0

        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]

            # Try to extract embedded images first
            image_list = page.get_images(full=True)

            for img_index, img_info in enumerate(image_list):
                xref = img_info[0]
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]
                image_size = len(image_bytes)

                # Keep track of the largest image
                if image_size > best_image_size:
                    best_image_size = image_size
                    # Convert to base64 with proper data URI
                    img_ext = base_image["ext"]
                    mime_type = f"image/{img_ext}" if img_ext != "jpg" else "image/jpeg"
                    best_image = f"data:{mime_type};base64,{base64.b64encode(image_bytes).decode('utf-8')}"

            # If no embedded images, render the page as an image
            if not image_list:
                # Render page at 2x resolution for better quality
                mat = fitz.Matrix(2, 2)
                pix = page.get_pixmap(matrix=mat)
                image_bytes = pix.tobytes("png")
                image_size = len(image_bytes)

                if image_size > best_image_size:
                    best_image_size = image_size
                    best_image = f"data:image/png;base64,{base64.b64encode(image_bytes).decode('utf-8')}"

        pdf_document.close()

        if best_image:
            return {
                "success": True,
                "image_base64": best_image,
                "image_size": best_image_size,
            }
        else:
            return {"success": False, "error": "No images found in PDF"}

    except Exception as e:
        print(f"Error extracting images from PDF: {str(e)}")
        import traceback

        traceback.print_exc()
        return {"success": False, "error": str(e)}


async def parse_pdf_smart(
    file_base64: str, prefer_image: bool = False
) -> Dict[str, Any]:
    """
    Smart PDF parser that detects content type and extracts accordingly.

    Args:
        file_base64: Base64 encoded PDF file
        prefer_image: If True, try to extract images first (for flowchart mode)

    Returns:
        Dictionary with either text content or image data
    """
    try:
        # Decode base64 content
        file_content = base64.b64decode(
            file_base64.split(",")[1] if "," in file_base64 else file_base64
        )

        # Open PDF with pymupdf to check content
        pdf_document = fitz.open(stream=file_content, filetype="pdf")

        has_images = False
        has_significant_text = False
        total_text_length = 0
        total_image_area = 0

        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]

            # Check for images
            image_list = page.get_images(full=True)
            if image_list:
                has_images = True
                for img_info in image_list:
                    # Estimate image area from image dimensions
                    xref = img_info[0]
                    try:
                        base_image = pdf_document.extract_image(xref)
                        total_image_area += len(base_image["image"])
                    except:
                        pass

            # Check for text
            text = page.get_text()
            total_text_length += len(text.strip())

        pdf_document.close()

        # Determine content type based on preferences and content analysis
        if prefer_image:
            # In flowchart mode, always try to get images
            image_result = await parse_pdf_for_images(file_base64)
            if image_result["success"]:
                return {
                    "success": True,
                    "content_type": "image",
                    "content": image_result["image_base64"],
                }

        # If we have significant text (more than 50 chars), extract as text
        if total_text_length > 50:
            has_significant_text = True

        # Decision logic
        if prefer_image and has_images:
            image_result = await parse_pdf_for_images(file_base64)
            if image_result["success"]:
                return {
                    "success": True,
                    "content_type": "image",
                    "content": image_result["image_base64"],
                }

        if has_significant_text:
            # Extract text
            file_stream = io.BytesIO(file_content)
            text, metadata = await parse_pdf(file_stream)
            return {
                "success": True,
                "content_type": "text",
                "content": text,
                "metadata": metadata,
            }

        # Fallback: try to render page as image
        image_result = await parse_pdf_for_images(file_base64)
        if image_result["success"]:
            return {
                "success": True,
                "content_type": "image",
                "content": image_result["image_base64"],
            }

        return {"success": False, "error": "Could not extract content from PDF"}

    except Exception as e:
        print(f"Error in smart PDF parsing: {str(e)}")
        import traceback

        traceback.print_exc()
        return {"success": False, "error": str(e)}
